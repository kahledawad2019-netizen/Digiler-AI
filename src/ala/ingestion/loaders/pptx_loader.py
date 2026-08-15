"""PowerPoint loader (.pptx) via python-pptx.

One slide -> a HEADING (title) + body blocks (bulleted text -> LIST, tables ->
TABLE), all anchored to the slide number, with the slide title as section_path.
This matches the doc-type-aware chunking plan (slide = parent).
"""

from __future__ import annotations

from pathlib import Path

from ala.core.enums import DocType, ExtractionMethod
from ala.fabric.content import Anchor, BlockType
from ala.ingestion.loaders.base import BaseLoader, BlockSpec


class PptxLoader(BaseLoader):
    name = "pptx"
    extensions = (".pptx", ".ppsx")   # .ppsx = PowerPoint Show; same OOXML format
    doc_types = (DocType.LECTURE_SLIDES.value,)
    extraction_method = ExtractionMethod.PPTX
    requires = ("pptx",)

    def _parse(self, path: Path) -> list[BlockSpec]:
        prs = self._open_presentation(path)
        specs: list[BlockSpec] = []

        for idx, slide in enumerate(prs.slides, start=1):
            anchor = Anchor(slide=idx)
            title = self._slide_title(slide)
            section = [title] if title else [f"Slide {idx}"]
            title_id = slide.shapes.title.shape_id if slide.shapes.title is not None else None
            if title:
                specs.append(BlockSpec(text=title, block_type=BlockType.HEADING,
                                       section_path=section, anchor=anchor, meta={"slide": idx}))

            for shape in slide.shapes:
                if shape.has_table:
                    specs.append(BlockSpec(text=self._table_text(shape.table),
                                           block_type=BlockType.TABLE,
                                           section_path=section, anchor=anchor))
                    continue
                if not shape.has_text_frame:
                    continue
                if title_id is not None and shape.shape_id == title_id:
                    continue  # already emitted as heading
                bullets = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                if not bullets:
                    continue
                btype = BlockType.LIST if len(bullets) > 1 else BlockType.PARAGRAPH
                specs.append(BlockSpec(text="\n".join(bullets), block_type=btype,
                                       section_path=section, anchor=anchor, meta={"slide": idx}))
        return specs

    @staticmethod
    def _open_presentation(path: Path):
        """Open a .pptx or .ppsx presentation.

        python-pptx rejects .ppsx because a PowerPoint Show declares the OOXML
        "slideshow" main content-type instead of "presentation". We rewrite that
        one content-type in an in-memory copy of the package (the parts are
        otherwise identical), then load normally.
        """
        from pptx import Presentation

        if path.suffix.lower() != ".ppsx":
            return Presentation(str(path))

        import io
        import zipfile

        buffer = io.BytesIO()
        with zipfile.ZipFile(str(path)) as src, \
                zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as out:
            for item in src.namelist():
                data = src.read(item)
                if item == "[Content_Types].xml":
                    data = data.replace(
                        b"presentationml.slideshow.main+xml",
                        b"presentationml.presentation.main+xml",
                    )
                out.writestr(item, data)
        buffer.seek(0)
        return Presentation(buffer)

    @staticmethod
    def _slide_title(slide) -> str | None:
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                return slide.shapes.title.text.strip()
        except (AttributeError, ValueError):
            pass
        return None

    @staticmethod
    def _table_text(table) -> str:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(rows)
