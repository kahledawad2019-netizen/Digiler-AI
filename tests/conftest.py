"""Shared pytest fixtures."""

from __future__ import annotations

import pytest

from ala.catalog.database import Database
from ala.catalog.repository import KnowledgeCatalog
from ala.config.settings import load_settings
from ala.core.clock import FixedClock
from ala.metadata.schema import FileInfo, ResourceMetadata


@pytest.fixture(scope="session")
def settings():
    return load_settings()


@pytest.fixture
def clock():
    return FixedClock("2026-07-08T00:00:00+00:00")


@pytest.fixture
def mem_catalog(clock):
    """A fresh in-memory catalog per test (fast, isolated)."""
    cat = KnowledgeCatalog(Database(":memory:"), clock=clock)
    cat.initialize()
    yield cat
    cat.close()


@pytest.fixture
def make_meta():
    """Factory for a valid ResourceMetadata without touching disk."""

    def _make(resource_id="technical.dmv.w03.constraints", **overrides) -> ResourceMetadata:
        base = dict(
            resource_id=resource_id,
            title="Constraints and Relationships",
            track="technical",
            course="dmv",
            module="w03",
            file=FileInfo(
                file_name="source.pdf",
                file_path="knowledge_base/raw/x/source.pdf",
                file_size=1234,
                sha256="a" * 64,
            ),
        )
        base.update(overrides)
        return ResourceMetadata(**base)

    return _make


@pytest.fixture
def sample_file(tmp_path):
    """A small on-disk file to register."""
    p = tmp_path / "lecture.txt"
    p.write_text("Foreign keys enforce referential integrity.\n", encoding="utf-8")
    return p


@pytest.fixture
def registry_factory(settings, mem_catalog):
    """Build a ResourceRegistry whose project_root is a chosen directory.

    Pointing project_root at a test's tmp_path makes registered file paths
    relative to that root, which is what change detection needs.
    """
    from ala.registry.registry import ResourceRegistry

    def _make(project_root):
        scoped = settings.model_copy(update={"project_root": project_root})
        return ResourceRegistry(scoped, mem_catalog)

    return _make


@pytest.fixture
def ingest_pipeline_factory(settings, mem_catalog):
    """Build an IngestionPipeline rooted at a chosen dir, backed by the in-memory catalog."""
    from ala.ingestion import IngestionPipeline
    from ala.registry.registry import ResourceRegistry

    def _make(project_root):
        scoped = settings.model_copy(update={"project_root": project_root})
        registry = ResourceRegistry(scoped, mem_catalog)
        return IngestionPipeline.default(scoped, registry=registry)

    return _make


# --- sample-file generators for each supported document type ---------------- #
@pytest.fixture
def make_docx():
    def _make(path):
        import docx

        d = docx.Document()
        d.add_heading("Constraints and Relationships", level=1)
        d.add_paragraph("A primary key uniquely identifies each row in a table.")
        d.add_heading("Foreign Keys", level=2)
        d.add_paragraph("A foreign key enforces referential integrity.", style="List Bullet")
        table = d.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Term"
        table.rows[0].cells[1].text = "Meaning"
        d.save(str(path))
        return path

    return _make


@pytest.fixture
def make_pptx():
    def _make(path):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "K-Nearest Neighbors"
        body = slide.placeholders[1].text_frame
        body.text = "Distance-based classifier"
        body.add_paragraph().text = "Choose k neighbours"
        prs.save(str(path))
        return path

    return _make


@pytest.fixture
def make_notebook():
    def _make(path):
        import nbformat as nbf

        nb = nbf.v4.new_notebook()
        nb.cells = [
            nbf.v4.new_markdown_cell("# Decision Trees\n\nA tree splits on features."),
            nbf.v4.new_code_cell("from sklearn.tree import DecisionTreeClassifier"),
        ]
        nbf.write(nb, str(path))
        return path

    return _make
